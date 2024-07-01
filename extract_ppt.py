import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from io import BytesIO

ppt_path = "investor-presentation-jefferies-2016.pptx" #filepath
presentation = Presentation(ppt_path)

images_dir = "ppt_images"
text_dir = "ppt_text"
os.makedirs(images_dir, exist_ok=True)
os.makedirs(text_dir, exist_ok=True)


def save_image(image, image_path):
    try:
        with open(image_path, 'wb') as f:
            f.write(image.blob)
    except Exception as e:
        print(f"Failed to save image at {image_path}: {e}")


def save_text(text, text_path):
    try:
        with open(text_path, 'w', encoding='utf-8') as f:
            f.write(text)
    except Exception as e:
        print(f"Failed to save text at {text_path}: {e}")

def extract_images_from_shape(shape, slide_index):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = shape.image
            image_filename = f"slide_{slide_index}_image_{shape.shape_id}.png"
            image_path = os.path.join(images_dir, image_filename)
            save_image(image, image_path)
        except Exception as e:
            print(f"Failed to process image on slide {slide_index}, shape {shape.shape_id}: {e}")
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            extract_images_from_shape(sub_shape, slide_index)


def extract_and_save_ppt_content():
    try:
        for i, slide in enumerate(presentation.slides, start=1):
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text += shape.text + '\n'
                extract_images_from_shape(shape, i)

            text_filename = f"slide_{i}.txt"
            text_path = os.path.join(text_dir, text_filename)
            save_text(slide_text, text_path)

        print("Extraction complete. Check the ppt_images and ppt_text folders.")
    except Exception as e:
        print(f"Failed to extract content from PPT: {e}")


extract_and_save_ppt_content()
